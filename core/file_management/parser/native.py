import os
import json
import logging
from typing import List, Optional, Dict, Any, TYPE_CHECKING
from urllib.parse import urlparse
from dotenv import load_dotenv

load_dotenv()

from .base import AbstractParser
from framework.singleton_decorator import singleton

if TYPE_CHECKING:
    from config.core.file_management.parser.native import NativeParserConfig

logger = logging.getLogger(__name__)


@singleton
class NativeParser(AbstractParser):
    """
    Multi-format document parser supporting PDF, DOCX, Excel, PowerPoint, HTML, and images.

    This class provides a unified parsing interface for multiple document formats,
    automatically routing files to appropriate specialized parsing functions based on file extension.
    Supports both local files and remote URLs with configurable output formats.

    Supported file types:
    - DOCX: Native Word document parsing with image extraction
    - Excel: XLSX/CSV parsing with table structure preservation
    - PowerPoint: PPTX slide content and layout extraction
    - HTML: Web page parsing with content extraction
    - Images: JPG/PNG (requires DotsOCR integration)
    - PDF: (requires DotsOCR integration)
    """

    def __init__(self, config: "NativeParserConfig"):
        """Initialize NativeParser"""
        super().__init__(config)

    async def parse_file(
        self,
        file_data: bytes,
        filename: str,
        **kwargs: Any
    ) -> List[Dict[str, Any]]:
        """Parse a file of any supported type from binary data"""

        # Get output directory from environment variable
        output_dir = os.getenv('NATIVE_PARSER_OUTPUT_DIR', './test_output/native')
        output_dir = os.path.abspath(output_dir)
        os.makedirs(output_dir, exist_ok=True)

        # Extract file extension and validate
        base_filename, file_ext = os.path.splitext(filename)
        file_ext = file_ext.lower()

        if file_ext not in self.get_supported_extensions():
            supported = ', '.join(self.get_supported_extensions())
            raise ValueError(f"Unsupported file type '{file_ext}'. Supported types: {supported}")

        # Route to appropriate parser method
        try:
            if file_ext == '.docx':
                return self._parse_docx(file_data, filename, output_dir, **kwargs)
            elif file_ext in ['.xlsx', '.xls', '.csv']:
                return self._parse_excel(file_data, filename, output_dir, **kwargs)
            elif file_ext == '.pptx':
                return self._parse_ppt(file_data, filename, output_dir, **kwargs)
            elif file_ext == '.html':
                return self._parse_html_content(file_data.decode('utf-8'), filename, base_filename, output_dir)
            elif file_ext == '.txt':
                return self._parse_txt(file_data, filename, output_dir, **kwargs)
            elif file_ext == '.md':
                return self._parse_md(file_data, filename, output_dir, **kwargs)
            else:
                raise ValueError(f"File type '{file_ext}' is listed as supported but no handler exists")

        except Exception as e:
            logger.error(f"Failed to parse {filename}: {str(e)}")
            raise RuntimeError(f"Failed to parse {filename}: {str(e)}")

    def get_supported_extensions(self) -> List[str]:
        """Get all supported file extensions"""
        return ['.docx', '.xlsx', '.xls', '.csv', '.pptx', '.html', '.txt', '.md']

    def _parse_docx(self, file_data: bytes, filename: str, output_dir: str, **kwargs) -> List[Dict[str, Any]]:
        """Parse DOCX file from binary data and return structured results"""
        try:
            import io
            from docx import Document

            base_filename = os.path.splitext(filename)[0]
            save_dir = os.path.join(output_dir, base_filename)
            os.makedirs(save_dir, exist_ok=True)

            logger.info(f"Parsing DOCX: {base_filename}")

            # Parse DOCX content from binary data
            doc = Document(io.BytesIO(file_data))

            # Extract text content
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)

            # Extract tables
            tables_data = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables_data.append(table_data)

            # Save results
            content = {
                'text': '\n'.join(full_text),
                'tables': tables_data,
                'metadata': {
                    'filename': base_filename,
                    'paragraphs_count': len([p for p in doc.paragraphs if p.text.strip()]),
                    'tables_count': len(tables_data)
                }
            }

            # Save as JSON
            json_path = os.path.join(save_dir, f"{base_filename}.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(content, f, ensure_ascii=False, indent=2)

            # Save as Markdown
            md_content = self._convert_docx_to_markdown(content)
            md_path = os.path.join(save_dir, f"{base_filename}.md")
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)

            result = {
                'filename': filename,
                'page_no': 0,
                'content_type': 'docx',
                'output_paths': {
                    'json': json_path,
                    'markdown': md_path
                },
                'metadata': content['metadata']
            }

            return [result]

        except Exception as e:
            logger.error(f"DOCX parsing failed: {str(e)}")
            raise

    def _parse_excel(self, file_data: bytes, filename: str, output_dir: str, **kwargs) -> List[Dict[str, Any]]:
        """Parse Excel file from binary data and return structured results"""
        try:
            import io
            import pandas as pd
            import chardet

            base_filename = os.path.splitext(filename)[0]
            file_ext = os.path.splitext(filename)[1].lower()
            save_dir = os.path.join(output_dir, base_filename)
            os.makedirs(save_dir, exist_ok=True)

            logger.info(f"Parsing Excel: {base_filename}, file size: {len(file_data)} bytes")

            # Validate file data
            if not file_data or len(file_data) == 0:
                raise ValueError(f"File data is empty for {filename}")

            # Check if file data looks like a valid Excel file (should start with PK for .xlsx)
            if file_ext == '.xlsx' and not file_data.startswith(b'PK'):
                logger.warning(f"File {filename} does not appear to be a valid .xlsx file (missing PK header)")
                logger.debug(f"First 100 bytes: {file_data[:100]}")

            # Read all sheets from binary data
            if file_ext == '.csv':
                # Detect encoding for CSV
                encoding_result = chardet.detect(file_data[:10000])
                encoding = encoding_result['encoding'] or "utf-8"
                csv_content = file_data.decode(encoding)
                sheets_data = {'Sheet1': pd.read_csv(io.StringIO(csv_content))}
            else:
                # Explicitly specify engine to avoid pandas auto-detection issues
                # .xlsx files use openpyxl, .xls files use xlrd
                if file_ext == '.xlsx':
                    engine = 'openpyxl'
                elif file_ext == '.xls':
                    # Check if xlrd is available for .xls files
                    try:
                        import xlrd
                        engine = 'xlrd'
                    except ImportError:
                        raise RuntimeError(
                            f"Cannot parse .xls file '{filename}': xlrd library is not installed. "
                            "Please install dependencies with: uv sync"
                        )
                else:
                    # Should not reach here based on supported extensions check
                    engine = None

                logger.debug(f"Using pandas engine: {engine} for {file_ext} file")
                sheets_data = pd.read_excel(io.BytesIO(file_data), sheet_name=None, engine=engine)

            all_content = []
            sheet_results = []

            for sheet_name, df in sheets_data.items():
                # Convert to structured data
                sheet_content = {
                    'sheet_name': sheet_name,
                    'data': df.to_dict('records'),
                    'columns': df.columns.tolist(),
                    'shape': df.shape,
                    'metadata': {
                        'rows': len(df),
                        'columns': len(df.columns),
                        'empty_cells': int(df.isnull().sum().sum())
                    }
                }
                all_content.append(sheet_content)

                # Save individual sheet as CSV
                csv_path = os.path.join(save_dir, f"{base_filename}_{sheet_name}.csv")
                df.to_csv(csv_path, index=False, encoding='utf-8')

                # Save as JSON
                json_path = os.path.join(save_dir, f"{base_filename}_{sheet_name}.json")
                with open(json_path, 'w', encoding='utf-8') as f:
                    json.dump(sheet_content, f, ensure_ascii=False, indent=2)

                sheet_results.append({
                    'filename': filename,
                    'page_no': len(sheet_results),
                    'content_type': 'excel_sheet',
                    'sheet_name': sheet_name,
                    'output_paths': {
                        'json': json_path,
                        'csv': csv_path
                    },
                    'metadata': sheet_content['metadata']
                })

            # Save combined results
            combined_json = os.path.join(save_dir, f"{base_filename}_combined.json")
            with open(combined_json, 'w', encoding='utf-8') as f:
                json.dump(all_content, f, ensure_ascii=False, indent=2)

            return sheet_results

        except Exception as e:
            logger.error(f"Excel parsing failed: {str(e)}")
            raise

    def _parse_ppt(self, file_data: bytes, filename: str, output_dir: str, **kwargs) -> List[Dict[str, Any]]:
        """Parse PowerPoint file from binary data and return structured results"""
        try:
            import io
            from pptx import Presentation

            base_filename = os.path.splitext(filename)[0]
            save_dir = os.path.join(output_dir, base_filename)
            os.makedirs(save_dir, exist_ok=True)

            logger.info(f"Parsing PowerPoint: {base_filename}")

            prs = Presentation(io.BytesIO(file_data))
            slides_data = []
            results = []

            for i, slide in enumerate(prs.slides):
                slide_content = {
                    'slide_number': i + 1,
                    'title': '',
                    'text_content': [],
                    'notes': ''
                }

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if not slide_content['title'] and hasattr(shape, 'text_frame'):
                            slide_content['title'] = shape.text.strip()
                        else:
                            slide_content['text_content'].append(shape.text.strip())

                # Extract notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_content['notes'] = slide.notes_slide.notes_text_frame.text.strip()

                slides_data.append(slide_content)

                # Save individual slide
                slide_json = os.path.join(save_dir, f"{base_filename}_slide_{i+1}.json")
                with open(slide_json, 'w', encoding='utf-8') as f:
                    json.dump(slide_content, f, ensure_ascii=False, indent=2)

                # Convert to markdown
                md_content = self._convert_slide_to_markdown(slide_content)
                slide_md = os.path.join(save_dir, f"{base_filename}_slide_{i+1}.md")
                with open(slide_md, 'w', encoding='utf-8') as f:
                    f.write(md_content)

                results.append({
                    'filename': filename,
                    'page_no': i,
                    'content_type': 'ppt_slide',
                    'slide_number': i + 1,
                    'output_paths': {
                        'json': slide_json,
                        'markdown': slide_md
                    },
                    'metadata': {
                        'title': slide_content['title'],
                        'text_blocks': len(slide_content['text_content']),
                        'has_notes': bool(slide_content['notes'])
                    }
                })

            # Save combined presentation
            combined_json = os.path.join(save_dir, f"{base_filename}_combined.json")
            with open(combined_json, 'w', encoding='utf-8') as f:
                json.dump(slides_data, f, ensure_ascii=False, indent=2)

            return results

        except Exception as e:
            logger.error(f"PowerPoint parsing failed: {str(e)}")
            raise

    def _parse_html_content(self, html_content: str, filename: str, base_filename: str, output_dir: str) -> List[Dict[str, Any]]:
        """Parse HTML content and return structured results"""
        try:
            from bs4 import BeautifulSoup

            save_dir = os.path.join(output_dir, base_filename)
            os.makedirs(save_dir, exist_ok=True)

            logger.info(f"Parsing HTML: {base_filename}")

            # Parse HTML
            soup = BeautifulSoup(html_content, 'html.parser')

            # Extract structured content
            content = {
                'title': soup.title.string if soup.title else '',
                'headings': [],
                'paragraphs': [],
                'links': [],
                'images': [],
                'tables': [],
                'metadata': {
                    'filename': base_filename
                }
            }

            # Extract headings
            for tag in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                for heading in soup.find_all(tag):
                    if heading.get_text().strip():
                        content['headings'].append({
                            'level': tag,
                            'text': heading.get_text().strip()
                        })

            # Extract paragraphs
            for p in soup.find_all('p'):
                text = p.get_text().strip()
                if text:
                    content['paragraphs'].append(text)

            # Extract links
            for a in soup.find_all('a', href=True):
                content['links'].append({
                    'text': a.get_text().strip(),
                    'href': a['href']
                })

            # Extract images
            for img in soup.find_all('img'):
                content['images'].append({
                    'src': img.get('src', ''),
                    'alt': img.get('alt', ''),
                    'title': img.get('title', '')
                })

            # Extract tables
            for table in soup.find_all('table'):
                table_data = []
                for row in table.find_all('tr'):
                    row_data = [cell.get_text().strip() for cell in row.find_all(['td', 'th'])]
                    if row_data:
                        table_data.append(row_data)
                if table_data:
                    content['tables'].append(table_data)

            # Save results
            json_path = os.path.join(save_dir, f"{base_filename}.json")
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(content, f, ensure_ascii=False, indent=2)

            # Convert to markdown
            md_content = self._convert_html_to_markdown(content)
            md_path = os.path.join(save_dir, f"{base_filename}.md")
            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)

            result = {
                'filename': filename,
                'page_no': 0,
                'content_type': 'html',
                'output_paths': {
                    'json': json_path,
                    'markdown': md_path
                },
                'metadata': {
                    'title': content['title'],
                    'headings_count': len(content['headings']),
                    'paragraphs_count': len(content['paragraphs']),
                    'links_count': len(content['links']),
                    'images_count': len(content['images']),
                    'tables_count': len(content['tables'])
                }
            }

            return [result]

        except Exception as e:
            logger.error(f"HTML parsing failed: {str(e)}")
            raise

    # ==================== PRIVATE UTILITY METHODS ====================

    def _convert_docx_to_markdown(self, content: dict) -> str:
        """Convert DOCX content to Markdown"""
        md_lines = []

        # Add text content
        if content['text']:
            md_lines.append(content['text'])
            md_lines.append('')

        # Add tables
        for i, table in enumerate(content['tables']):
            md_lines.append(f"## Table {i+1}")
            md_lines.append('')

            if table and len(table) > 0:
                # Header row
                header = '| ' + ' | '.join(table[0]) + ' |'
                separator = '| ' + ' | '.join(['---'] * len(table[0])) + ' |'
                md_lines.append(header)
                md_lines.append(separator)

                # Data rows
                for row in table[1:]:
                    row_md = '| ' + ' | '.join(row) + ' |'
                    md_lines.append(row_md)

            md_lines.append('')

        return '\n'.join(md_lines)

    def _convert_slide_to_markdown(self, slide_content: dict) -> str:
        """Convert slide content to Markdown"""
        md_lines = []

        if slide_content['title']:
            md_lines.append(f"# {slide_content['title']}")
            md_lines.append('')

        for text in slide_content['text_content']:
            md_lines.append(text)
            md_lines.append('')

        if slide_content['notes']:
            md_lines.append("## Notes")
            md_lines.append(slide_content['notes'])

        return '\n'.join(md_lines)

    def _convert_html_to_markdown(self, content: dict) -> str:
        """Convert HTML content to Markdown"""
        md_lines = []

        if content['title']:
            md_lines.append(f"# {content['title']}")
            md_lines.append('')

        # Add headings and paragraphs in order they appear
        for heading in content['headings']:
            level = int(heading['level'][1])
            md_lines.append('#' * level + f" {heading['text']}")
            md_lines.append('')

        for paragraph in content['paragraphs']:
            md_lines.append(paragraph)
            md_lines.append('')

        # Add tables
        for i, table in enumerate(content['tables']):
            md_lines.append(f"## Table {i+1}")
            md_lines.append('')

            if table and len(table) > 0:
                # Header row
                header = '| ' + ' | '.join(table[0]) + ' |'
                separator = '| ' + ' | '.join(['---'] * len(table[0])) + ' |'
                md_lines.append(header)
                md_lines.append(separator)

                # Data rows
                for row in table[1:]:
                    row_md = '| ' + ' | '.join(row) + ' |'
                    md_lines.append(row_md)

            md_lines.append('')

        return '\n'.join(md_lines)

    def _parse_txt(self, file_data: bytes, filename: str, output_dir: str, **kwargs) -> List[Dict[str, Any]]:
        """Parse plain text file from binary data and return structured results"""
        try:
            # Decode text content
            text_content = file_data.decode('utf-8', errors='ignore')

            # Create output directory
            base_filename = os.path.splitext(filename)[0]
            save_dir = os.path.join(output_dir, base_filename)
            os.makedirs(save_dir, exist_ok=True)

            # Save as markdown (plain text is valid markdown)
            md_filename = f"{base_filename}.md"
            md_path = os.path.join(save_dir, md_filename)

            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(text_content)

            logger.info(f"TXT parsing complete: {filename} -> {md_path}")

            return [{
                "type": "text",
                "content": text_content,
                "metadata": {
                    "source_file": filename,
                    "output_file": md_path,
                    "format": "txt"
                }
            }]

        except Exception as e:
            logger.error(f"TXT parsing failed: {str(e)}")
            raise

    def _parse_md(self, file_data: bytes, filename: str, output_dir: str, **kwargs) -> List[Dict[str, Any]]:
        """
        Parse Markdown file from binary data - direct save without conversion.

        For .md files, we simply save them as-is since they're already in Markdown format.
        No conversion or processing is needed.
        """
        try:
            # Decode markdown content
            md_content = file_data.decode('utf-8', errors='ignore')

            # Create output directory
            base_filename = os.path.splitext(filename)[0]
            save_dir = os.path.join(output_dir, base_filename)
            os.makedirs(save_dir, exist_ok=True)

            # Save markdown file directly (no conversion needed)
            md_filename = f"{base_filename}.md"
            md_path = os.path.join(save_dir, md_filename)

            with open(md_path, 'w', encoding='utf-8') as f:
                f.write(md_content)

            logger.info(f"MD file saved directly: {filename} -> {md_path}")

            return [{
                "type": "markdown",
                "content": md_content,
                "metadata": {
                    "source_file": filename,
                    "output_file": md_path,
                    "format": "md"
                }
            }]

        except Exception as e:
            logger.error(f"MD file processing failed: {str(e)}")
            raise